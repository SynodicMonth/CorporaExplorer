import docx
import pdfplumber
import pptx


def pdf2text(path):
    with pdfplumber.open(path) as pdf:
        text = ''
        if len(pdf.pages) > 30:
            for page in pdf.pages[:30]:
                text += page.extract_text()
        else:
            for page in pdf.pages:
                text += page.extract_text()
        return text


def doc2text(path):
    doc = docx.Document(path)
    text = ''
    for para in doc.paragraphs:
        text += para.text
    return text


def txt2text(path):
    with open(path, 'r') as f:
        return f.read()


def ppt2text(path):
    ppt = pptx.Presentation(path)
    text = ''
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text
    return text


def get_text(path, file_type):
    if file_type == 'pdf':
        return pdf2text(path)
    elif file_type == 'docx' or file_type == 'doc':
        return doc2text(path)
    elif file_type in ['txt', 'py', 'c', 'cpp', 'java', 'html', 'css', 'js', 'php', 'sql', 'xml', 'json', 'md']:
        return txt2text(path)
    elif file_type == 'pptx' or file_type == 'ppt':
        return ppt2text(path)
    else:
        raise Exception('Unsupported file type')
